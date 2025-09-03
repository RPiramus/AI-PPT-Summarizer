from fastapi import FastAPI, File, UploadFile, HTTPException, status, APIRouter
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from dotenv import load_dotenv
from openai import OpenAI
from pydantic import BaseModel
from fastapi.responses import FileResponse
import os, subprocess
from .db import Base, engine
from .routers import auth

os.makedirs("converted", exist_ok=True)
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

Base.metadata.create_all(bind=engine)

app = FastAPI()

app.include_router(auth.router)

origins = [ "http://localhost:5173", "http://127.0.0.1:5173",]
app.add_middleware( CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"], 
)

frontend_path = os.path.join(os.path.dirname(__file__),"..", "..", "frontend")
app.mount("/static", StaticFiles(directory=frontend_path), name="static")
app.mount("/converted", StaticFiles(directory="converted"), name="converted")

app.include_router(auth.router, prefix="/api/auth", tags=["auth"])

api = APIRouter(prefix="/api")

@api.get("/health")
def health():
    return {"ok": True}

@app.post("/upload")
async def upload_ppt(file: UploadFile = File(...)):
    ext      = os.path.splitext(file.filename)[1].lower() 
    src_path = f"temp{ext}"
    contents = await file.read()
    with open(src_path, "wb") as f:
        f.write(contents)

    if ext == ".ppt":
        subprocess.run([
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "--headless", "--convert-to", "pptx",
            "--outdir", ".", src_path
        ], check=True)
        pptx_path = "temp.pptx"
    else:
        pptx_path = src_path

    prs = Presentation(pptx_path)
    extracted = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    prompt = (
        "You are an assistant summarizing PowerPoint slides into clean, structured study notes.\n"
        "Each slide is labeled. For each slide:\n\n"
        "1. Start with a slide topic (e.g., 'Topic 1: Chapter Title')\n"
        "2. Include an 'Overview' or 'Key Points' section\n"
        "3. Use bullet points with •\n"
        "4. Use **bold** for key terms or section headers\n"
        "5. If the slide contains math or formulas, use code blocks with LaTeX notation (like this):\n"
        "```\n"
        "\\text{Example Equation} = 2x + 3\n"
        "```\n"
        "6. Do not force equations or code blocks unless it’s appropriate\n"
        "7. Always include every slide, even if it has little content\n\n"
        "Here is the presentation content:\n\n"
        + extracted
    )

    resp = client.chat.completions.create(
        model="gpt-3.5-turbo-16k",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that summarizes PowerPoint presentations."},
            {"role": "user",   "content": prompt},
        ],
        temperature=0.7,
        max_tokens=4000,
    )
    summary = resp.choices[0].message.content

    subprocess.run([
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "--headless", "--convert-to", "pdf",
        "--outdir", "converted", pptx_path
    ], check=True)

    return {
        "message": "PPT uploaded successfully",
        "summary": summary,
        "pdf_url": "/converted/temp.pdf"
    }


class ChatRequest(BaseModel):
    summary: str
    question: str

@app.post("/chat")
async def chat_with_summary(request: ChatRequest):
    messages = [
        {"role": "system", "content": "You are a helpful assistant that answers questions based on a PPT summary."},
        {"role": "user",   "content": f"Here’s the summary:\n\n{request.summary}"},
        {"role": "user",   "content": request.question},
    ]
    resp = client.chat.completions.create(
        model="gpt-3.5-turbo-16k",
        messages=messages,
        temperature=0.7,
        max_tokens=4000,
    )
    return {"answer": resp.choices[0].message.content}

# source venv/bin/activate
# uvicorn app.main:app --reload